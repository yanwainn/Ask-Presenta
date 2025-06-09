"""
PowerPoint conversion functionality using organization template.
Converts content to PowerPoint slides using a corporate template.
"""
import os
import io
import base64
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image

def create_powerpoint_from_template(slides, extraction_result=None, template_path=None, output_path=None, template_stream=None, template_file=None):
    """
    Create a PowerPoint presentation from slides using organization template.
    
    Args:
        slides (list): List of HTMLSlide objects
        extraction_result: Content extraction result
        template_path (str): Path to the template PPTX file
        output_path (str): Path to save the presentation
        template_stream (BytesIO): BytesIO stream containing the template PPTX
        template_file (str): Alternative name for template_path
        
    Returns:
        BytesIO: BytesIO object containing the PowerPoint file
    """
    # Handle template_file parameter (alias for template_path)
    if template_file and not template_path:
        template_path = template_file
        
    # Check if template exists - either as a file or stream
    if template_stream:
        # Use the provided BytesIO stream
        prs = Presentation(template_stream)
    elif template_path and os.path.exists(template_path):
        # Load from file path
        prs = Presentation(template_path)
    else:
        raise FileNotFoundError(f"Template file not found: {template_path}")
    
    # Store the original slides from the template for reference
    template_slides = [slide for slide in prs.slides]
    
    # Clean the presentation - remove all but the first slide (usually the title slide)
    # We'll keep only the first slide as title and use it as a template
    while len(prs.slides) > 1:
        r_id = prs.slides._sldIdLst[-1].rId
        prs.part.drop_rel(r_id)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[-1])
    
    # Add the title to the first slide (simplified - just the title with no subtitle)
    if len(prs.slides) > 0:
        title_slide = prs.slides[0]
        
        # Find title placeholder
        title_placeholder = None
        subtitle_placeholder = None
        
        for shape in title_slide.shapes:
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 1:  # Title
                title_placeholder = shape
            elif hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 2:  # Subtitle
                subtitle_placeholder = shape
        
        # Add title if found
        if title_placeholder and hasattr(title_placeholder, 'text_frame'):
            if extraction_result and hasattr(extraction_result, 'document_title'):
                title_placeholder.text_frame.text = extraction_result.document_title
            else:
                title_placeholder.text_frame.text = "Presentation"
        
        # Clear subtitle if found (or add minimal subtitle if needed)
        if subtitle_placeholder and hasattr(subtitle_placeholder, 'text_frame'):
            subtitle_placeholder.text_frame.text = ""  # Leave subtitle empty as requested
    
    # Find content slide templates in the original template
    content_layout = None
    
    # Look for a content layout with title and content placeholders
    for slide_layout in prs.slide_layouts:
        has_title = False
        has_content = False
        
        for placeholder in slide_layout.placeholders:
            if placeholder.placeholder_format.type == 1:  # Title
                has_title = True
            if placeholder.placeholder_format.type == 7:  # Content
                has_content = True
        
        if has_title and has_content:
            content_layout = slide_layout
            break
    
    # If no suitable layout found, use the first layout after title slide
    if not content_layout and len(prs.slide_layouts) > 1:
        content_layout = prs.slide_layouts[1]
    
    # If still no layout, use blank layout
    if not content_layout:
        content_layout = prs.slide_layouts[6]  # Usually blank layout
    
    # Try to load company logo if exists (for blank layouts where we need to add manually)
    logo_path = None
    for p in [os.path.join(os.getcwd(), "logo.png"), 
              os.path.join(os.getcwd(), "static", "logo.png"),
              os.path.join(os.getcwd(), "assets", "logo.png")]:
        if os.path.exists(p):
            logo_path = p
            break
    
    # Process each slide
    for slide_obj in slides:
        # Create a new slide using the template layout
        slide = prs.slides.add_slide(content_layout)
        
        # Find title and content placeholders
        title_shape = None
        content_shape = None
        image_placeholder = None
        
        for shape in slide.placeholders:
            if hasattr(shape, 'placeholder_format'):
                # Common placeholder types:
                # 1: Title
                # 7: Content
                # 18: Picture
                if shape.placeholder_format.type == 1:  # Title
                    title_shape = shape
                elif shape.placeholder_format.type == 7:  # Content
                    content_shape = shape
                elif shape.placeholder_format.type in [18, 17, 19]:  # Picture or media placeholders
                    image_placeholder = shape
        
        # Add title if placeholder exists
        if title_shape:
            title_shape.text = slide_obj.title
        
        # Process bullet points from section_content
        bullet_points = []
        if slide_obj.section_content and "," in slide_obj.section_content:
            bullet_points = slide_obj.section_content.split(", ")
        else:
            # If no comma-separated content, try to extract bullet points from HTML
            bullet_pattern = r'<div class="bullet-text">([^<]+)</div>'
            bullet_matches = re.findall(bullet_pattern, slide_obj.html_content)
            if bullet_matches:
                bullet_points = bullet_matches
        
        # Filter out empty bullet points
        bullet_points = [p for p in bullet_points if p and len(p.strip()) > 0]
        
        # Add bullet points to content placeholder if it exists
        if content_shape and bullet_points:
            text_frame = content_shape.text_frame
            text_frame.clear()  # Clear any default text
            
            # Add each bullet point
            for i, point in enumerate(bullet_points):
                p = text_frame.add_paragraph()
                p.text = point
                p.level = 0  # Main bullet level
        
        # Extract and add image 
        base64_image_pattern = r'src="data:image\/jpeg;base64,([^"]+)"'
        match = re.search(base64_image_pattern, slide_obj.html_content)
        
        if match:
            # Extract base64 image data
            base64_data = match.group(1)
            image_data = base64.b64decode(base64_data)
            image_stream = io.BytesIO(image_data)
            
            # Try inserting image
            if image_placeholder:
                # Insert into image placeholder if available
                try:
                    image_placeholder.insert_picture(image_stream)
                except Exception as e:
                    print(f"Error inserting into placeholder: {e}")
                    # Fallback - add as regular shape
                    image_width = Inches(4)
                    image_left = prs.slide_width - image_width - Inches(0.5)
                    image_top = Inches(2)
                    slide.shapes.add_picture(image_stream, image_left, image_top, width=image_width)
            else:
                # Add as regular shape if no placeholder
                # Position depends on whether we have content
                if content_shape:
                    # If there's a content block, position image on the right
                    image_width = Inches(4)
                    image_left = prs.slide_width - image_width - Inches(0.5)
                    image_top = Inches(2)
                else:
                    # If no content, center the image
                    image_width = Inches(6)
                    image_left = (prs.slide_width - image_width) / 2
                    image_top = (prs.slide_height - Inches(5)) / 2
                
                # Add the image
                slide.shapes.add_picture(image_stream, image_left, image_top, width=image_width)
    
    # Save the presentation
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    
    # Also save to file if path is provided
    if output_path:
        prs.save(output_path)
    
    return output