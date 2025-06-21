"""
PowerPoint conversion functionality.
Convert HTML slides to PowerPoint format with styling similar to HTML slides.
"""
import os
import io
import base64
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image

# BOSCH Colors - matching the HTML slides
BOSCH_COLORS = {
    "primary": "#E20015",  # Bosch Red
    "black": "#000000",
    "white": "#FFFFFF",
    "dark_gray": "#333333",
    "medium_gray": "#7D7D7D",
    "light_gray": "#E5E5E5",
    "secondary_blue": "#007BC0",
    "secondary_green": "#00884A"
}

def hex_to_rgb(hex_color):
    """
    Convert hex color to RGB.
    
    Args:
        hex_color (str): Hex color code (e.g., "#FF0000")
        
    Returns:
        tuple: (R, G, B) values
    """
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def create_powerpoint_from_slides(slides, extraction_result=None, output_path=None):
    """
    Create a PowerPoint presentation from HTML slides with similar styling.
    
    Args:
        slides (list): List of HTMLSlide objects
        extraction_result: Content extraction result
        output_path (str): Path to save the presentation
        
    Returns:
        BytesIO: BytesIO object containing the PowerPoint file
    """
    # Create a new presentation
    prs = Presentation()
    
    # Set slide size to 16:9 aspect ratio
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # Create a custom slide layout without bullet points (we'll add our own content)
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    title_slide_layout = prs.slide_layouts[0]  # Title slide layout
    
    # Extract colors from BOSCH_COLORS
    primary_color = hex_to_rgb(BOSCH_COLORS["primary"])  # Bosch Red
    secondary_color = hex_to_rgb(BOSCH_COLORS["secondary_blue"])  # Bosch Blue
    text_color = hex_to_rgb(BOSCH_COLORS["dark_gray"])  # Dark Gray
    light_text_color = hex_to_rgb(BOSCH_COLORS["medium_gray"])  # Medium Gray
    
    # Add a title slide
    title_slide = prs.slides.add_slide(title_slide_layout)
    
    # Style the title slide
    title = title_slide.shapes.title
    title.text = extraction_result.document_title if extraction_result else "Presentation"
    
    # Style the title text on title slide
    title_text_frame = title.text_frame
    title_paragraph = title_text_frame.paragraphs[0]
    title_paragraph.alignment = PP_ALIGN.LEFT
    title_run = title_paragraph.runs[0]
    title_run.font.size = Pt(44)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(*text_color)
    
    # Add subtitle
    subtitle = extraction_result.summary if extraction_result else ""
    if subtitle and len(title_slide.placeholders) > 1:
        subtitle_placeholder = title_slide.placeholders[1]
        subtitle_placeholder.text = subtitle[:300] + "..." if len(subtitle) > 300 else subtitle
        
        # Style the subtitle
        for paragraph in subtitle_placeholder.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.size = Pt(20)
                run.font.color.rgb = RGBColor(*secondary_color)
    
    # Try to load company logo if exists
    logo_path = None
    for p in [os.path.join(os.getcwd(), "logo.png"), 
              os.path.join(os.getcwd(), "static", "logo.png"),
              os.path.join(os.getcwd(), "assets", "logo.png")]:
        if os.path.exists(p):
            logo_path = p
            break
    
    # Add logo to title slide if available
    if logo_path:
        logo_height = Inches(0.6)
        logo_left = prs.slide_width - Inches(1.2)
        logo_top = prs.slide_height - Inches(0.8)
        title_slide.shapes.add_picture(logo_path, logo_left, logo_top, height=logo_height)
    
    # Process each slide
    for slide_obj in slides:
        # Create a blank slide (we'll add custom elements)
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add a title box at the top
        title_left = Inches(0.5)
        title_top = Inches(0.5)
        title_width = prs.slide_width - Inches(1)
        title_height = Inches(1)
        
        title_shape = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_frame = title_shape.text_frame
        title_paragraph = title_frame.add_paragraph()
        title_paragraph.text = slide_obj.title
        title_paragraph.alignment = PP_ALIGN.LEFT
        
        # Style the title
        title_run = title_paragraph.runs[0]
        title_run.font.size = Pt(36)
        title_run.font.bold = True
        title_run.font.color.rgb = RGBColor(*text_color)
        
        # Add subtitle with document title
        subtitle_paragraph = title_frame.add_paragraph()
        subtitle_paragraph.text = extraction_result.document_title if extraction_result else ""
        subtitle_paragraph.alignment = PP_ALIGN.LEFT
        
        # Style the subtitle
        subtitle_run = subtitle_paragraph.runs[0]
        subtitle_run.font.size = Pt(20)
        subtitle_run.font.color.rgb = RGBColor(*secondary_color)
        
        # Process bullet points from section_content
        bullet_points = []
        if slide_obj.section_content and "," in slide_obj.section_content:
            bullet_points = slide_obj.section_content.split(", ")
        else:
            # If no comma-separated content, try to extract bullet points from HTML
            bullet_pattern = r'<div class="bullet-text">([^<]+)</div>'
            bullet_matches = re.findall(bullet_pattern, slide_obj.html_content)
            if bullet_matches:
                bullet_points = bullet_matches
        
        # Filter out empty bullet points
        bullet_points = [p for p in bullet_points if p and len(p.strip()) > 0]
        
        # Determine layout based on content - similar to the HTML slide logic
        # If few bullet points -> image focus, if many -> text focus
        is_image_focus = len(bullet_points) <= 2 and "image" in slide_obj.title.lower()
        is_text_focus = len(bullet_points) >= 5
        
        # Extract base64 image from the slide
        base64_image_pattern = r'src="data:image\/jpeg;base64,([^"]+)"'
        match = re.search(base64_image_pattern, slide_obj.html_content)
        image_data = None
        
        if match:
            # Extract base64 image data
            base64_data = match.group(1)
            image_data = base64.b64decode(base64_data)
        
        # Position content based on layout
        if is_image_focus:
            # Image focus layout - large image on top, bullets below
            if image_data:
                image_stream = io.BytesIO(image_data)
                image_width = Inches(10)
                image_left = (prs.slide_width - image_width) / 2  # Center
                image_top = Inches(1.8)
                slide.shapes.add_picture(image_stream, image_left, image_top, width=image_width)
            
            # Add bullet points below the image
            bullet_box_left = Inches(1)
            bullet_box_top = Inches(5.5)
            bullet_box_width = prs.slide_width - Inches(2)
            bullet_box_height = Inches(2.5)
            
            if bullet_points:
                bullet_box = slide.shapes.add_textbox(bullet_box_left, bullet_box_top, bullet_box_width, bullet_box_height)
                
                # Create styled bullet points
                bullet_frame = bullet_box.text_frame
                for i, point in enumerate(bullet_points):
                    p = bullet_frame.add_paragraph()
                    p.text = point
                    p.level = 0
                    
                    # Style the bullet point
                    run = p.runs[0]
                    run.font.size = Pt(20)
                    run.font.color.rgb = RGBColor(*text_color)
        else:
            # Standard layout - bullets on left, image on right
            # Or text focus - more space for text
            
            # Calculate sizes based on layout
            if is_text_focus:
                # Text focus - bullets get more space
                bullet_box_width = prs.slide_width * 0.7
                image_width = prs.slide_width * 0.25
            else:
                # Balanced layout
                bullet_box_width = prs.slide_width * 0.55
                image_width = prs.slide_width * 0.35
            
            # Add bullet points on the left
            bullet_box_left = Inches(0.5)
            bullet_box_top = Inches(2)
            bullet_box_height = prs.slide_height - Inches(3)
            
            if bullet_points:
                bullet_box = slide.shapes.add_textbox(bullet_box_left, bullet_box_top, bullet_box_width, bullet_box_height)
                
                # Create a styled container to match HTML slide look
                bullet_frame = bullet_box.text_frame
                
                for i, point in enumerate(bullet_points):
                    p = bullet_frame.add_paragraph()
                    
                    # Add bullet character (similar to HTML)
                    p.text = f"• {point}"
                    p.level = 0
                    
                    # Style the bullet point
                    run = p.runs[0]
                    run.font.size = Pt(18)
                    run.font.color.rgb = RGBColor(*text_color)
                    
                    # Add space between bullet points
                    if i < len(bullet_points) - 1:
                        p.space_after = Pt(12)
            
            # Add image on the right if available
            if image_data:
                image_stream = io.BytesIO(image_data)
                image_left = prs.slide_width - image_width - Inches(0.5)
                image_top = Inches(2)
                slide.shapes.add_picture(image_stream, image_left, image_top, width=image_width)
        
        # Add footer with document title and slide title
        footer_top = prs.slide_height - Inches(0.5)
        footer_left = Inches(0.5)
        footer_width = prs.slide_width - Inches(2)
        footer_height = Inches(0.4)
        
        footer_shape = slide.shapes.add_textbox(footer_left, footer_top, footer_width, footer_height)
        footer_frame = footer_shape.text_frame
        footer_paragraph = footer_frame.add_paragraph()
        footer_paragraph.text = f"{extraction_result.document_title if extraction_result else 'Presentation'} | {slide_obj.title}"
        footer_paragraph.alignment = PP_ALIGN.LEFT
        
        # Style the footer
        footer_run = footer_paragraph.runs[0]
        footer_run.font.size = Pt(10)
        footer_run.font.color.rgb = RGBColor(*light_text_color)
        
        # Add company logo if available
        if logo_path:
            logo_height = Inches(0.5)
            logo_left = prs.slide_width - Inches(1.2)
            logo_top = prs.slide_height - Inches(0.8)
            slide.shapes.add_picture(logo_path, logo_left, logo_top, height=logo_height)
    
    # Save the presentation
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    
    # Also save to file if path is provided
    if output_path:
        prs.save(output_path)
    
    return output